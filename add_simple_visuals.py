from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

print("Loading presentation...")

# Load the presentation
prs = Presentation(r'C:\Users\torre\Downloads\oklo-consulting-presentation\OKLO - From Reactors to Revenue.pptx')

# Oklo colors
OKLO_TEAL = RGBColor(72, 201, 176)
OKLO_NAVY = RGBColor(0, 51, 102)
WHITE = RGBColor(255, 255, 255)

print(f"Total slides: {len(prs.slides)}")

# First, remove ALL shapes from slides 2 and 5 except title/text boxes with actual content
def clean_slide(slide, slide_num):
    print(f"\nCleaning Slide {slide_num}...")
    shapes_to_remove = []

    for shape in slide.shapes:
        # Keep only shapes that have meaningful text content
        if hasattr(shape, 'text'):
            text = shape.text.strip()
            # Keep if it has substantial text (more than just whitespace)
            if len(text) > 5:
                continue

        # Remove everything else (empty text boxes, shapes, icons, etc.)
        shapes_to_remove.append(shape)

    # Remove shapes
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)

    print(f"  Removed {len(shapes_to_remove)} shapes")

# Clean both slides
clean_slide(prs.slides[1], 2)
clean_slide(prs.slides[4], 5)

# =====================================================
# SLIDE 2: Simple number visual
# Just a large "80+" with clean subtitle
# =====================================================

print("\nAdding SIMPLE visual to Slide 2...")

slide2 = prs.slides[1]

# Single centered visual - just the number
number_box = slide2.shapes.add_textbox(
    Inches(2.5), Inches(2.5), Inches(5.0), Inches(2.0)
)
number_frame = number_box.text_frame
number_frame.word_wrap = True

# Big number
number_para = number_frame.paragraphs[0]
number_para.text = "80+"
number_para.font.size = Pt(120)
number_para.font.bold = True
number_para.font.color.rgb = OKLO_TEAL
number_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_para = number_frame.add_paragraph()
subtitle_para.text = "startups need regulatory help"
subtitle_para.font.size = Pt(24)
subtitle_para.font.color.rgb = WHITE
subtitle_para.alignment = PP_ALIGN.CENTER
subtitle_para.space_before = Pt(10)

print("  DONE - Added clean centered number")

# =====================================================
# SLIDE 5: Simple arrow visual
# Just 2020 → 2025 with simple text
# =====================================================

print("\nAdding SIMPLE visual to Slide 5...")

slide5 = prs.slides[4]

# Create simple left-to-right flow
# Left box: 2020
left_box = slide5.shapes.add_textbox(
    Inches(1.5), Inches(2.5), Inches(2.5), Inches(1.8)
)
left_frame = left_box.text_frame
left_frame.word_wrap = True

left_year = left_frame.paragraphs[0]
left_year.text = "2020"
left_year.font.size = Pt(48)
left_year.font.bold = True
left_year.font.color.rgb = WHITE
left_year.alignment = PP_ALIGN.CENTER

left_desc = left_frame.add_paragraph()
left_desc.text = "First NRC\nApplication"
left_desc.font.size = Pt(18)
left_desc.font.color.rgb = WHITE
left_desc.alignment = PP_ALIGN.CENTER
left_desc.space_before = Pt(10)

# Arrow in middle
arrow_box = slide5.shapes.add_textbox(
    Inches(4.2), Inches(3.0), Inches(1.6), Inches(1.0)
)
arrow_frame = arrow_box.text_frame
arrow_para = arrow_frame.paragraphs[0]
arrow_para.text = "→"
arrow_para.font.size = Pt(72)
arrow_para.font.color.rgb = OKLO_TEAL
arrow_para.alignment = PP_ALIGN.CENTER

# Right box: 2025
right_box = slide5.shapes.add_textbox(
    Inches(6.0), Inches(2.5), Inches(2.5), Inches(1.8)
)
right_frame = right_box.text_frame
right_frame.word_wrap = True

right_year = right_frame.paragraphs[0]
right_year.text = "2025"
right_year.font.size = Pt(48)
right_year.font.bold = True
right_year.font.color.rgb = OKLO_TEAL
right_year.alignment = PP_ALIGN.CENTER

right_desc = right_frame.add_paragraph()
right_desc.text = "$140M\nExpertise"
right_desc.font.size = Pt(18)
right_desc.font.color.rgb = OKLO_TEAL
right_desc.alignment = PP_ALIGN.CENTER
right_desc.space_before = Pt(10)

print("  DONE - Added simple 2020 to 2025 flow")

# Save the presentation
output_file = r'C:\Users\torre\Downloads\oklo-consulting-presentation\OKLO - From Reactors to Revenue.pptx'
prs.save(output_file)

print(f"\nSUCCESS! Clean, simple visuals added: {output_file}")
print("\nSlide 2: Just 80+ in large text")
print("Slide 5: Simple 2020 to 2025 flow")
print("\nNo clutter. No overlapping. Clean and professional.")
