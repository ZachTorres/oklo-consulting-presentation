from pptx import Presentation

# Load the Google Slides PowerPoint
prs = Presentation('OKLO - From Reactors to Revenue.pptx')

print("=" * 80)
print("OKLO - From Reactors to Revenue - Slide Analysis")
print("=" * 80)
print(f"\nTotal Slides: {len(prs.slides)}\n")

# Analyze each slide
for slide_num, slide in enumerate(prs.slides, start=1):
    print(f"\n{'=' * 80}")
    print(f"SLIDE {slide_num}")
    print(f"{'=' * 80}")

    # Get all text from the slide
    all_text = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    all_text.append(text)

        # Check for charts
        if shape.shape_type == 3:  # Chart
            print("\n[CHART DETECTED]")
            try:
                chart = shape.chart
                print(f"Chart Type: {chart.chart_type}")
                if chart.has_title:
                    print(f"Chart Title: {chart.chart_title.text_frame.text}")
            except:
                pass

        # Check for tables
        if shape.shape_type == 19:  # Table
            print("\n[TABLE DETECTED]")
            try:
                table = shape.table
                print(f"Table: {len(table.rows)} rows x {len(table.columns)} columns")
            except:
                pass

    # Print all text content
    if all_text:
        print("\nText Content:")
        for i, text in enumerate(all_text, 1):
            print(f"  {i}. {text}")
    else:
        print("\n[No text content]")

print("\n" + "=" * 80)
print("Analysis Complete!")
print("=" * 80)
