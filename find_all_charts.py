from pptx import Presentation

# Load the PowerPoint
prs = Presentation('OKLO - From Reactors to Revenue.pptx')

print(f"Total Slides: {len(prs.slides)}\n")

# Look through all slides
for slide_idx, slide in enumerate(prs.slides):
    slide_num = slide_idx + 1

    # Get slide title if exists
    title_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and len(text) < 100:
                title_text = text
                break

    # Look for charts
    has_chart = False
    for shape in slide.shapes:
        if hasattr(shape, "has_chart") and shape.has_chart:
            has_chart = True
            print(f"SLIDE {slide_num}: {title_text}")
            print(f"  >>> CHART FOUND <<<")
            chart = shape.chart
            print(f"  Chart type: {chart.chart_type}")
            print(f"  Number of series: {len(chart.series)}")
            for series_idx, series in enumerate(chart.series):
                print(f"    Series {series_idx + 1}: {series.name}")
            print()

    if not has_chart and title_text:
        print(f"SLIDE {slide_num}: {title_text}")
