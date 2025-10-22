from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Load the PowerPoint
prs = Presentation('OKLO - From Reactors to Revenue.pptx')

slide = prs.slides[7]  # Slide 8 (0-indexed as 7)

print("=" * 80)
print("SLIDE 8 - All Shapes Analysis")
print("=" * 80)

for shape_idx, shape in enumerate(slide.shapes):
    print(f"\nShape {shape_idx}:")
    print(f"  Name: {shape.name}")
    print(f"  Type: {shape.shape_type} ({MSO_SHAPE_TYPE(shape.shape_type) if shape.shape_type in [e.value for e in MSO_SHAPE_TYPE] else 'UNKNOWN'})")

    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            print(f"  Text: {text[:100]}")

    # Check if it's a chart
    if shape.shape_type == MSO_SHAPE_TYPE.CHART or shape.shape_type == 3:
        print("  >>> THIS IS A CHART! <<<")

    # Check if it's a graphic
    if hasattr(shape, 'chart'):
        print("  >>> HAS CHART ATTRIBUTE! <<<")
