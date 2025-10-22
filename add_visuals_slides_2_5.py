from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import math

print("Loading presentation...")

# Load the presentation
prs = Presentation(r'C:\Users\torre\Downloads\oklo-consulting-presentation\OKLO - From Reactors to Revenue.pptx')

# Oklo colors
OKLO_TEAL = RGBColor(72, 201, 176)
OKLO_NAVY = RGBColor(0, 51, 102)
OKLO_LIGHT_BLUE = RGBColor(93, 173, 226)
OKLO_DARK_GRAY = RGBColor(51, 51, 51)
WHITE = RGBColor(255, 255, 255)
OKLO_ORANGE = RGBColor(255, 102, 0)

print(f"Total slides: {len(prs.slides)}")

# =====================================================
# SLIDE 2: "80+ New nuclear startups need help"
# Visual: Icon grid showing 80 small atoms with some highlighted
# =====================================================

print("\nEnhancing Slide 2: Icon Grid Visualization...")

slide2 = prs.slides[1]  # Index 1 = Slide 2

# Create 80 small atom/reactor icons in a grid
# Grid: 10 columns x 8 rows = 80 icons
grid_cols = 10
grid_rows = 8
total_icons = 80

# Positioning
start_x = Inches(0.5)
start_y = Inches(2.8)
icon_size = Inches(0.35)
spacing_x = Inches(0.95)
spacing_y = Inches(0.55)

print(f"  Creating {total_icons} reactor icons in {grid_cols}x{grid_rows} grid...")

# Highlight indices (representing potential Oklo clients)
highlight_indices = [5, 12, 23, 34, 47, 56, 63, 71, 78]  # 9 highlighted

icon_count = 0
for row in range(grid_rows):
    for col in range(grid_cols):
        if icon_count >= total_icons:
            break

        # Calculate position
        x = start_x + (col * spacing_x)
        y = start_y + (row * spacing_y)

        # Create circular shape (atom/reactor icon)
        shape = slide2.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, y, icon_size, icon_size
        )

        # Determine if this icon should be highlighted
        is_highlighted = icon_count in highlight_indices

        if is_highlighted:
            # Highlighted icons - Oklo teal with glow
            shape.fill.solid()
            shape.fill.fore_color.rgb = OKLO_TEAL
            shape.line.color.rgb = OKLO_TEAL
            shape.line.width = Pt(2)
            # Add glow effect
            shape.shadow.inherit = False
        else:
            # Regular icons - light gray/white outline
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
            shape.line.color.rgb = RGBColor(150, 150, 150)
            shape.line.width = Pt(1)

        icon_count += 1

# Add legend/caption
legend_box = slide2.shapes.add_textbox(
    Inches(7.5), Inches(3.0), Inches(2.2), Inches(1.5)
)
legend_frame = legend_box.text_frame
legend_frame.word_wrap = True

# Legend item 1
p1 = legend_frame.paragraphs[0]
p1.text = "● 80+ Startups"
p1.font.size = Pt(14)
p1.font.color.rgb = RGBColor(150, 150, 150)
p1.font.bold = True

# Legend item 2
p2 = legend_frame.add_paragraph()
p2.text = "● Target Clients"
p2.font.size = Pt(14)
p2.font.color.rgb = OKLO_TEAL
p2.font.bold = True
p2.space_before = Pt(8)

# Add count visualization
count_box = slide2.shapes.add_textbox(
    Inches(7.5), Inches(4.8), Inches(2.2), Inches(1.0)
)
count_frame = count_box.text_frame
count_para = count_frame.paragraphs[0]
count_para.text = "~$4B+"
count_para.font.size = Pt(36)
count_para.font.bold = True
count_para.font.color.rgb = OKLO_TEAL
count_para.alignment = PP_ALIGN.LEFT

count_para2 = count_frame.add_paragraph()
count_para2.text = "in total funding"
count_para2.font.size = Pt(12)
count_para2.font.color.rgb = WHITE

print("  DONE - Added 80 reactor icons with 9 highlighted targets")

# =====================================================
# SLIDE 5: "Oklo has already done this"
# Visual: Before/After timeline with transformation
# =====================================================

print("\nEnhancing Slide 5: Transformation Timeline...")

slide5 = prs.slides[4]  # Index 4 = Slide 5

# Timeline positions
timeline_y = Inches(3.2)
timeline_length = Inches(8.5)
timeline_start_x = Inches(0.75)

# Draw timeline base line
timeline = slide5.shapes.add_connector(
    1,  # Straight connector
    timeline_start_x, timeline_y,
    timeline_start_x + timeline_length, timeline_y
)
timeline.line.color.rgb = OKLO_TEAL
timeline.line.width = Pt(3)

# Milestone 1: 2020 Application
milestone1_x = timeline_start_x + Inches(0.5)
milestone1_y = timeline_y - Inches(0.6)

# Icon circle
m1_circle = slide5.shapes.add_shape(
    MSO_SHAPE.OVAL,
    milestone1_x, milestone1_y, Inches(0.8), Inches(0.8)
)
m1_circle.fill.solid()
m1_circle.fill.fore_color.rgb = OKLO_LIGHT_BLUE
m1_circle.line.color.rgb = WHITE
m1_circle.line.width = Pt(3)

# Add text inside circle
m1_textbox = slide5.shapes.add_textbox(
    milestone1_x, milestone1_y + Inches(0.15), Inches(0.8), Inches(0.5)
)
m1_frame = m1_textbox.text_frame
m1_para = m1_frame.paragraphs[0]
m1_para.text = "2020"
m1_para.font.size = Pt(16)
m1_para.font.bold = True
m1_para.font.color.rgb = WHITE
m1_para.alignment = PP_ALIGN.CENTER

# Label below
m1_label = slide5.shapes.add_textbox(
    milestone1_x - Inches(0.5), timeline_y + Inches(0.3), Inches(1.8), Inches(0.8)
)
m1_label_frame = m1_label.text_frame
m1_label_para = m1_label_frame.paragraphs[0]
m1_label_para.text = "First NRC\nApplication"
m1_label_para.font.size = Pt(14)
m1_label_para.font.color.rgb = WHITE
m1_label_para.alignment = PP_ALIGN.CENTER
m1_label_para.font.bold = True

# Milestone 2: 2022 Rejection (Pivot Point)
milestone2_x = timeline_start_x + Inches(4.0)
milestone2_y = timeline_y - Inches(0.6)

# Icon - Red X transforms to lightbulb
m2_circle = slide5.shapes.add_shape(
    MSO_SHAPE.OVAL,
    milestone2_x, milestone2_y, Inches(0.8), Inches(0.8)
)
m2_circle.fill.solid()
m2_circle.fill.fore_color.rgb = OKLO_ORANGE
m2_circle.line.color.rgb = WHITE
m2_circle.line.width = Pt(3)

# Text inside
m2_textbox = slide5.shapes.add_textbox(
    milestone2_x, milestone2_y + Inches(0.15), Inches(0.8), Inches(0.5)
)
m2_frame = m2_textbox.text_frame
m2_para = m2_frame.paragraphs[0]
m2_para.text = "2022"
m2_para.font.size = Pt(16)
m2_para.font.bold = True
m2_para.font.color.rgb = WHITE
m2_para.alignment = PP_ALIGN.CENTER

# Label below
m2_label = slide5.shapes.add_textbox(
    milestone2_x - Inches(0.5), timeline_y + Inches(0.3), Inches(1.8), Inches(0.8)
)
m2_label_frame = m2_label.text_frame
m2_label_para = m2_label_frame.paragraphs[0]
m2_label_para.text = "NRC Rejection\n→ Learning"
m2_label_para.font.size = Pt(14)
m2_label_para.font.color.rgb = WHITE
m2_label_para.alignment = PP_ALIGN.CENTER
m2_label_para.font.bold = True

# Curved arrow showing transformation
arrow = slide5.shapes.add_connector(
    2,  # Curved connector
    milestone2_x + Inches(0.4), milestone2_y + Inches(0.4),
    milestone2_x + Inches(2.5), milestone2_y + Inches(0.4)
)
arrow.line.color.rgb = OKLO_TEAL
arrow.line.width = Pt(4)

# Milestone 3: 2025 Asset
milestone3_x = timeline_start_x + Inches(7.5)
milestone3_y = timeline_y - Inches(0.6)

# Icon - Trophy/Star
m3_circle = slide5.shapes.add_shape(
    MSO_SHAPE.OVAL,
    milestone3_x, milestone3_y, Inches(0.8), Inches(0.8)
)
m3_circle.fill.solid()
m3_circle.fill.fore_color.rgb = OKLO_TEAL
m3_circle.line.color.rgb = WHITE
m3_circle.line.width = Pt(3)

# Add glow effect
m3_circle.shadow.inherit = False

# Text inside
m3_textbox = slide5.shapes.add_textbox(
    milestone3_x, milestone3_y + Inches(0.15), Inches(0.8), Inches(0.5)
)
m3_frame = m3_textbox.text_frame
m3_para = m3_frame.paragraphs[0]
m3_para.text = "2025"
m3_para.font.size = Pt(16)
m3_para.font.bold = True
m3_para.font.color.rgb = WHITE
m3_para.alignment = PP_ALIGN.CENTER

# Label below
m3_label = slide5.shapes.add_textbox(
    milestone3_x - Inches(0.7), timeline_y + Inches(0.3), Inches(2.2), Inches(1.0)
)
m3_label_frame = m3_label.text_frame
m3_label_para = m3_label_frame.paragraphs[0]
m3_label_para.text = "$140M\nKnowledge Asset"
m3_label_para.font.size = Pt(16)
m3_label_para.font.color.rgb = OKLO_TEAL
m3_label_para.alignment = PP_ALIGN.CENTER
m3_label_para.font.bold = True

# Add checklist of learnings above timeline
checklist_x = Inches(1.0)
checklist_y = Inches(1.3)

checklist_items = [
    "✓ NRC Requirements Mastered",
    "✓ Regulatory Process Mapped",
    "✓ Common Mistakes Identified",
    "✓ Winning Strategy Developed"
]

for idx, item in enumerate(checklist_items):
    check_box = slide5.shapes.add_textbox(
        checklist_x + (idx * Inches(2.2)), checklist_y, Inches(2.0), Inches(0.4)
    )
    check_frame = check_box.text_frame
    check_para = check_frame.paragraphs[0]
    check_para.text = item
    check_para.font.size = Pt(11)
    check_para.font.color.rgb = OKLO_TEAL
    check_para.font.bold = True

print("  DONE - Added transformation timeline with 3 milestones")
print("  DONE - Added learning checklist")

# Save the presentation
output_file = r'C:\Users\torre\Downloads\oklo-consulting-presentation\OKLO - From Reactors to Revenue.pptx'
prs.save(output_file)

print(f"\nSUCCESS! Enhanced presentation saved: {output_file}")
print("\nSlide 2 Enhancements:")
print("  - 80 reactor icons in 10x8 grid")
print("  - 9 icons highlighted in Oklo teal (target clients)")
print("  - Legend showing startup categories")
print("  - $4B+ funding stat callout")
print("\nSlide 5 Enhancements:")
print("  - Timeline with 3 key milestones (2020, 2022, 2025)")
print("  - Visual transformation from rejection to asset")
print("  - Learning checklist above timeline")
print("  - Color-coded progression (blue → orange → teal)")
