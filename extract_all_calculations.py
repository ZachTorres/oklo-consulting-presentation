from pptx import Presentation
import sys

# Set UTF-8 encoding for output
sys.stdout.reconfigure(encoding='utf-8')

# Load the presentation
prs = Presentation('OKLO - From Reactors to Revenue.pptx')

print("=" * 100)
print("COMPLETE NUMERICAL DATA EXTRACTION - OKLO PRESENTATION")
print("=" * 100)
print(f"\nTotal Slides: {len(prs.slides)}\n")

all_numbers = []

# Analyze each slide
for slide_num, slide in enumerate(prs.slides, start=1):
    print(f"\n{'=' * 100}")
    print(f"SLIDE {slide_num}")
    print(f"{'=' * 100}")

    slide_numbers = []

    # Get all text from the slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    print(f"  • {text}")

                    # Extract numbers
                    if any(char.isdigit() for char in text):
                        slide_numbers.append({
                            'slide': slide_num,
                            'text': text
                        })

        # Check for charts
        if shape.shape_type == 3:  # Chart
            print("\n  [CHART DETECTED]")
            try:
                chart = shape.chart
                print(f"  Chart Type: {chart.chart_type}")

                # Try to extract chart data
                if hasattr(chart, 'plots'):
                    for plot in chart.plots:
                        print(f"  Plot categories: {plot.categories if hasattr(plot, 'categories') else 'N/A'}")
                        if hasattr(plot, 'series'):
                            for series in plot.series:
                                print(f"  Series: {series.name if hasattr(series, 'name') else 'Unnamed'}")
                                if hasattr(series, 'values'):
                                    print(f"    Values: {series.values}")

            except Exception as e:
                print(f"  Error reading chart: {e}")

        # Check for tables
        if shape.shape_type == 19:  # Table
            print("\n  [TABLE DETECTED]")
            try:
                table = shape.table
                print(f"  Table: {len(table.rows)} rows x {len(table.columns)} columns")
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    print(f"    Row {row_idx + 1}: {' | '.join(row_data)}")
            except Exception as e:
                print(f"  Error reading table: {e}")

    all_numbers.extend(slide_numbers)

print("\n" + "=" * 100)
print("SUMMARY OF ALL NUMERICAL DATA")
print("=" * 100)

for item in all_numbers:
    print(f"Slide {item['slide']}: {item['text']}")

print("\n" + "=" * 100)
print("Analysis Complete!")
print("=" * 100)
