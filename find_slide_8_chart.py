from pptx import Presentation
from pptx.chart.data import CategoryChartData

# Load the PowerPoint
prs = Presentation('OKLO - From Reactors to Revenue.pptx')

print(f"Total Slides: {len(prs.slides)}\n")

# Look at slide 8 (index 7)
if len(prs.slides) >= 8:
    slide = prs.slides[7]  # Slide 8 (0-indexed)

    print("=" * 80)
    print("SLIDE 8 - Looking for charts")
    print("=" * 80)

    for shape_idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "has_chart") and shape.has_chart:
            print(f"\nFound chart at shape index {shape_idx}")
            chart = shape.chart
            print(f"Chart type: {chart.chart_type}")

            # Print series information
            print(f"\nNumber of series: {len(chart.series)}")
            for series_idx, series in enumerate(chart.series):
                print(f"\nSeries {series_idx + 1}: {series.name}")
                print(f"  Current fill color: {series.format.fill.type}")

else:
    print("Presentation has fewer than 8 slides")
