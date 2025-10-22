from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor

# Load the PowerPoint
prs = Presentation('OKLO - From Reactors to Revenue.pptx')

# Get slide 8 (index 7)
slide = prs.slides[7]

print("=" * 80)
print("Updating Slide 8 - Revenue Projection Chart")
print("=" * 80)

# Remove the old picture chart (Shape 2 based on our analysis)
# We need to find and remove it
shape_to_remove = None
for shape in slide.shapes:
    if shape.shape_type == 13:  # PICTURE type
        shape_to_remove = shape
        break

if shape_to_remove:
    print(f"Removing old chart image: {shape_to_remove.name}")
    sp = shape_to_remove.element
    sp.getparent().remove(sp)

# Create chart data for the revenue projection
chart_data = CategoryChartData()
chart_data.categories = ['2026', '2027', '2028', '2029', '2030']

# Add three series with the revenue data
chart_data.add_series('Licensing & Regulatory', (3, 15, 38, 65, 75))
chart_data.add_series('Safety & Simulation', (1, 6, 20, 35, 42))
chart_data.add_series('Fuel Cycle', (0.5, 2, 9, 18, 23))

# Add the chart to the slide
# Position it in the center-bottom area
x = Inches(1.5)
y = Inches(2.5)
cx = Inches(7)
cy = Inches(4)

chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.AREA_STACKED, x, y, cx, cy, chart_data
)

chart = chart_shape.chart

# Configure chart appearance
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.TOP
chart.legend.include_in_layout = False

# Apply blue color scheme to each series
# Series 1: Licensing & Regulatory - Dark Navy Blue
series1 = chart.series[0]
fill1 = series1.format.fill
fill1.solid()
fill1.fore_color.rgb = RGBColor(0, 51, 102)  # Oklo Navy #003366

# Series 2: Safety & Simulation - Medium Blue
series2 = chart.series[1]
fill2 = series2.format.fill
fill2.solid()
fill2.fore_color.rgb = RGBColor(41, 128, 185)  # Medium Blue #2980B9

# Series 3: Fuel Cycle - Light Blue
series3 = chart.series[2]
fill3 = series3.format.fill
fill3.solid()
fill3.fore_color.rgb = RGBColor(52, 152, 219)  # Light Blue #3498DB

print("\nApplied new blue color scheme:")
print("  Series 1 (Licensing & Regulatory): Dark Navy #003366 - RGB(0, 51, 102)")
print("  Series 2 (Safety & Simulation): Medium Blue #2980B9 - RGB(41, 128, 185)")
print("  Series 3 (Fuel Cycle): Light Blue #3498DB - RGB(52, 152, 219)")

# Save the updated presentation
output_filename = 'OKLO - From Reactors to Revenue.pptx'
prs.save(output_filename)

print(f"\n✓ Updated presentation saved as: {output_filename}")
print("=" * 80)
